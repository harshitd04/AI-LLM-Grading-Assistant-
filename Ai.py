import streamlit as st
import os
import tempfile
from pathlib import Path
import zipfile
from docx import Document
from docx.shared import Inches
import PyPDF2
import pandas as pd
from pptx import Presentation
import io
from langchain_openai import ChatOpenAI
from langchain_anthropic import ChatAnthropic
from langchain.schema import HumanMessage
import base64

# Page config
st.set_page_config(
    page_title="AI Project Grading Assistant",
    page_icon="📚",
    layout="wide"
)

# Title and description
st.title("📚 AI Project Grading Assistant")
st.markdown("Upload student project files and get AI-generated feedback using your own API key.")

# Sidebar for API configuration
with st.sidebar:
    st.header("🔑 API Configuration")
    
    # Provider selection
    provider = st.selectbox(
        "Select AI Provider:",
        ["ChatGPT", "Claude"],
        index=0,  # Default to ChatGPT
        help="Choose between OpenAI ChatGPT or Anthropic Claude"
    )
    
    # API key input with dynamic label
    api_key_label = "Enter your OpenAI API Key:" if provider == "ChatGPT" else "Enter your Anthropic API Key:"
    api_key = st.text_input(
        api_key_label,
        type="password",
        help="Your API key will not be stored and is only used for this session"
    )
    
    # Model selection based on provider
    if provider == "ChatGPT":
        model_name = st.selectbox(
            "Select ChatGPT Model:",
            ["gpt-4o", "gpt-4", "gpt-4-turbo-preview", "gpt-3.5-turbo"],
            index=0,  # Default to gpt-4o (ChatGPT 4.0)
            help="Choose the ChatGPT model for generating feedback"
        )
    else:  # Claude
        # Claude model display names and their corresponding API names
        claude_model_options = {
            "Claude Sonnet 4": "claude-3-5-sonnet-20241022",
            "Claude Opus 4": "claude-3-opus-20240229", 
            "Claude Haiku 4": "claude-3-5-haiku-20241022"
        }
        
        selected_display_name = st.selectbox(
            "Select Claude Model:",
            list(claude_model_options.keys()),
            index=0,  # Default to Claude Sonnet 4
            help="Choose the Claude model for generating feedback"
        )
        
        # Get the actual API model name
        model_name = claude_model_options[selected_display_name]
    
    st.markdown("---")
    st.markdown("**Supported File Types:**")
    st.markdown("- PDF files")
    st.markdown("- PowerPoint (.pptx)")
    st.markdown("- Word documents (.docx)")

def extract_text_from_pdf(file):
    """Extract text from PDF file"""
    try:
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        return f"Error reading PDF: {str(e)}"

def extract_text_from_docx(file):
    """Extract text from DOCX file"""
    try:
        doc = Document(file)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    except Exception as e:
        return f"Error reading DOCX: {str(e)}"

def extract_text_from_pptx(file):
    """Extract text from PPTX file"""
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"

def process_file(uploaded_file):
    """Process uploaded file and extract content"""
    file_extension = Path(uploaded_file.name).suffix.lower()
    
    if file_extension == '.pdf':
        return extract_text_from_pdf(uploaded_file)
    elif file_extension == '.docx':
        return extract_text_from_docx(uploaded_file)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(uploaded_file)
    else:
        return "Unsupported file format"

def generate_feedback(content, student_name, file_name, api_key, model_name, provider):
    """Generate feedback using OpenAI or Anthropic API via LangChain"""
    try:
        # Initialize the appropriate model based on provider
        if provider == "ChatGPT":
            llm = ChatOpenAI(
                api_key=api_key,
                model=model_name,
                temperature=0.7
            )
        else:  # Claude
            llm = ChatAnthropic(
                anthropic_api_key=api_key,
                model=model_name,
                temperature=0.7
            )
        
        # Create the grading prompt
        prompt = f"""
You are an experienced academic instructor grading student projects. 

Student Name: {student_name}
Project File: {file_name}

Project Content:
{content}

Please provide detailed feedback on this student project following this structure:

1. **Overall Assessment** (Grade: A/B/C/D/F)
   - Brief summary of the project quality
   
2. **Strengths**
   - What the student did well
   - Specific examples from their work
   
3. **Areas for Improvement**
   - Specific areas that need work
   - Constructive suggestions
   
4. **Technical Quality**
   - Content organization and structure
   - Clarity of presentation
   
5. **Recommendations**
   - Specific next steps for improvement
   - Resources or techniques to explore

Please be constructive, specific, and encouraging while maintaining academic standards.
"""
        
        # Generate response using the updated method
        response = llm.invoke([HumanMessage(content=prompt)])
        return response.content
        
    except Exception as e:
        return f"Error generating feedback: {str(e)}"

def create_word_document(feedback, student_name, file_name):
    """Create a Word document with the feedback"""
    doc = Document()
    
    # Add title
    title = doc.add_heading(f'Project Feedback: {student_name}', 0)
    
    # Add project details
    doc.add_heading('Project Details', level=1)
    doc.add_paragraph(f'Student: {student_name}')
    doc.add_paragraph(f'File: {file_name}')
    doc.add_paragraph(f'Generated: {st.session_state.get("timestamp", "N/A")}')
    
    # Add feedback content
    doc.add_heading('AI-Generated Feedback', level=1)
    doc.add_paragraph(feedback)
    
    # Add note for teacher
    doc.add_heading('Teacher Notes', level=1)
    doc.add_paragraph('Please review and edit the above feedback as needed before sharing with the student.')
    
    # Save to bytes
    doc_bytes = io.BytesIO()
    doc.save(doc_bytes)
    doc_bytes.seek(0)
    
    return doc_bytes

# Main application
def main():
    if not api_key:
        provider_name = "OpenAI" if provider == "ChatGPT" else "Anthropic"
        st.warning(f"⚠️ Please enter your {provider_name} API key in the sidebar to continue.")
        st.stop()
    
    # File upload section
    st.header("📁 Upload Student Project")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        uploaded_files = st.file_uploader(
            "Choose project files",
            type=['pdf', 'docx', 'pptx'],
            accept_multiple_files=True,
            help="You can upload multiple files for a single student"
        )
    
    with col2:
        student_name = st.text_input(
            "Student Name:",
            placeholder="Enter student name"
        )
    
    if uploaded_files and student_name:
        st.success(f"📁 {len(uploaded_files)} file(s) uploaded for {student_name}")
        
        # Display uploaded files
        st.subheader("Uploaded Files:")
        for file in uploaded_files:
            st.write(f"• {file.name} ({file.size} bytes)")
        
        # Process button
        if st.button("🤖 Generate Feedback", type="primary"):
            with st.spinner("Processing files and generating feedback..."):
                
                # Extract content from all files
                all_content = ""
                for file in uploaded_files:
                    st.write(f"Processing: {file.name}")
                    content = process_file(file)
                    all_content += f"\n\n--- Content from {file.name} ---\n{content}"
                
                # Generate feedback
                feedback = generate_feedback(
                    all_content, 
                    student_name, 
                    ", ".join([f.name for f in uploaded_files]),
                    api_key,
                    model_name,
                    provider
                )
                
                # Store in session state
                st.session_state['feedback'] = feedback
                st.session_state['student_name'] = student_name
                st.session_state['file_names'] = [f.name for f in uploaded_files]
                st.session_state['timestamp'] = st.session_state.get('timestamp', 
                    f"{st.session_state.get('timestamp', pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S'))}")
                
                st.success("✅ Feedback generated successfully!")
    
    # Display feedback and download option
    if 'feedback' in st.session_state:
        st.header("📝 Generated Feedback")
        
        # Display feedback
        with st.expander("View Feedback", expanded=True):
            st.markdown(st.session_state['feedback'])
        
        # Create and offer download
        st.subheader("💾 Download Feedback")
        
        # Create Word document
        doc_bytes = create_word_document(
            st.session_state['feedback'],
            st.session_state['student_name'],
            ", ".join(st.session_state['file_names'])
        )
        
        # Download button
        st.download_button(
            label="📄 Download Feedback (Word Document)",
            data=doc_bytes.getvalue(),
            file_name=f"{st.session_state['student_name']}_feedback.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        
        # Clear button
        if st.button("🗑️ Clear and Start New"):
            for key in ['feedback', 'student_name', 'file_names', 'timestamp']:
                if key in st.session_state:
                    del st.session_state[key]
            st.rerun()

# Footer
st.markdown("---")
st.markdown("**Academy AI Grading Assistant** - Developed for efficient project evaluation")

if __name__ == "__main__":
    main()